#!/usr/bin/env python3
"""Convert text blobs in signal/ trajectories to multimodal files (PDF/DOCX/HTML/XLSX/PPTX/MP3/PNG).

Pipeline: blob (.blob) -> [Gemini Flash parse] -> structured JSON -> [local libs] -> multimodal files

Usage:
    python scripts/convert_multimodal.py --parallel 4
    python scripts/convert_multimodal.py --trajectory p4_structured_analyst_T-20
    python scripts/convert_multimodal.py --skip-existing
    python scripts/convert_multimodal.py --dry
"""

import argparse
import json
import logging
import os
import re
import sys
import time
import traceback
import warnings
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


# Suppress noisy font glyph warnings from fpdf2
logging.getLogger("fpdf2").setLevel(logging.ERROR)
warnings.filterwarnings("ignore", message=".*DeprecationWarning.*uni.*")

# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class BlobInfo:
    hash: str
    blob_path: Path
    original_path: str
    size_bytes: int
    text: str = ""


@dataclass
class ConversionResult:
    hash: str
    original_path: str
    modalities: dict = field(default_factory=dict)   # modality -> {path, size}
    errors: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SIGNAL_DIR = Path(__file__).resolve().parent.parent / "signal"
MAX_TEXT_CHARS = 30_000

# CJK font search paths (macOS)
CJK_FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Songti.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/Library/Fonts/Arial Unicode.ttf",
]

GEMINI_SCHEMA = {
    "type": "object",
    "properties": {
        "title": {"type": "string"},
        "language": {"type": "string", "enum": ["zh", "en"]},
        "modalities": {
            "type": "array",
            "items": {
                "type": "string",
                "enum": ["pdf", "docx", "xlsx", "pptx", "mp3", "png"],
            },
        },

        # PDF — condensed formal document
        "pdf_text": {"type": "string"},

        # DOCX — full detailed working document
        "docx_text": {"type": "string"},

        # PPTX — presentation slides with bullet points
        "pptx_slides": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "slide_title": {"type": "string"},
                    "bullets": {"type": "string"},
                },
                "required": ["slide_title", "bullets"],
            },
        },

        # MP3 — conversational narration script
        "mp3_script": {"type": "string"},

        # PNG — chart data or key quote
        "has_chart": {"type": "boolean"},
        "chart_type": {"type": "string", "enum": ["bar", "pie", "line", "none"]},
        "chart_labels": {"type": "array", "items": {"type": "string"}},
        "chart_values": {"type": "array", "items": {"type": "number"}},
        "chart_title": {"type": "string"},
        "key_quote": {"type": "string"},

        # XLSX — table data
        "has_table": {"type": "boolean"},
        "table_headers": {"type": "array", "items": {"type": "string"}},
        "table_rows": {
            "type": "array",
            "items": {"type": "array", "items": {"type": "string"}},
        },
    },
    "required": [
        "title", "language", "modalities",
        "pdf_text", "docx_text", "pptx_slides", "mp3_script",
        "has_chart", "chart_type", "chart_labels", "chart_values", "chart_title",
        "key_quote",
        "has_table", "table_headers", "table_rows",
    ],
}

GEMINI_PROMPT = """\
You are a document analyst. Given a text document, decide which 1-3 non-text file formats
best represent this content, then produce the content for each chosen format.

IMPORTANT: Be very selective. Most documents need only 1 or 2 formats. Maximum 3.
- ~40% of documents: 1 format (the single best fit).
- ~40% of documents: 2 formats (only if they serve clearly different purposes).
- ~20% of documents: 3 formats (only for rich documents with mixed content types, e.g. text + table + chart).
Do NOT pick 3 by default. Ask yourself: "Does this document REALLY need a third format?"

Decision guide for `modalities`:
- "pdf": good for formal reports, briefs, executive summaries. Use `pdf_text`.
- "docx": good for detailed working documents, full analyses, long-form content. Use `docx_text`.
- "pptx": good for presentations, multi-point arguments, structured takeaways. Use `pptx_slides`.
- "xlsx": ONLY if the document contains real tabular data. Use `table_headers`/`table_rows`.
- "mp3": good for conversational briefings, short narrations. Use `mp3_script`.
- "png": ONLY if there is numeric data for a chart OR a strong key quote. Use chart fields or `key_quote`.

Content rules:
- `pdf_text`: Condensed formal summary (≤500 words). Standalone 1-page brief.
- `docx_text`: Complete detailed document with full context (≤3000 words).
- `pptx_slides`: 3-8 slides, each with `slide_title` and `bullets` (concise phrases).
- `mp3_script`: Spoken-language narration (≤200 words), as if briefing a colleague.
- If you pick multiple formats, each must carry DIFFERENT content — no duplication.
- Only populate the content fields for formats you actually chose in `modalities`.
  For unchosen formats, use empty string or empty array.

Document:
---
{text}
---"""


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def find_cjk_font() -> Optional[str]:
    for p in CJK_FONT_CANDIDATES:
        if os.path.exists(p):
            return p
    return None


CJK_FONT_PATH = find_cjk_font()


def has_cjk(text: str) -> bool:
    return bool(re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', text))


def load_content_blobs(traj_dir: Path) -> list[BlobInfo]:
    """Read manifest.json, filter out fs_snapshot entries, load blob text."""
    manifest_path = traj_dir / "media" / "manifest.json"
    if not manifest_path.exists():
        return []

    with open(manifest_path) as f:
        manifest = json.load(f)

    entries = manifest.get("entries", {})
    blobs = []
    blobs_dir = traj_dir / "media" / "blobs"

    for hash_id, entry in entries.items():
        if entry.get("type") != "blob":
            continue
        orig = entry.get("original_path", "")
        # Skip fs_snapshot blobs
        if "fs_snapshot" in orig:
            continue

        blob_path = blobs_dir / f"{hash_id}.blob"
        if not blob_path.exists():
            continue

        try:
            text = blob_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            continue

        if not text.strip():
            continue

        blobs.append(BlobInfo(
            hash=hash_id,
            blob_path=blob_path,
            original_path=orig,
            size_bytes=entry.get("size_bytes", 0),
            text=text,
        ))

    return blobs


# ---------------------------------------------------------------------------
# Gemini client
# ---------------------------------------------------------------------------

def create_gemini_client():
    """Create Gemini client using google-genai SDK."""
    from google import genai

    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        # Try loading from .env
        env_path = Path(__file__).resolve().parent.parent / ".env"
        if env_path.exists():
            for line in env_path.read_text().splitlines():
                if line.startswith("GEMINI_API_KEY="):
                    api_key = line.split("=", 1)[1].strip().strip('"').strip("'")
                    break

    if not api_key:
        raise RuntimeError("GEMINI_API_KEY not found in env or .env file")

    client = genai.Client(api_key=api_key)
    return client


def analyze_blob(client, blob: BlobInfo) -> dict:
    """Call Gemini Flash to analyze a blob and return structured JSON."""
    from google.genai import types

    truncated = blob.text[:MAX_TEXT_CHARS]
    prompt = GEMINI_PROMPT.format(text=truncated)

    for attempt in range(3):
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.2,
                    response_mime_type="application/json",
                    response_schema=GEMINI_SCHEMA,
                ),
            )
            result = json.loads(response.text)
            return result
        except Exception as e:
            if attempt < 2 and ("429" in str(e) or "Resource" in str(e)):
                time.sleep(2 ** attempt + 1)
                continue
            print(f"\r  [WARN] Gemini failed for {blob.hash[:12]}: {e}")
            return None


# ---------------------------------------------------------------------------
# Generators
# ---------------------------------------------------------------------------

def generate_html(blob: BlobInfo, analysis: Optional[dict], out_dir: Path) -> Optional[Path]:
    """Convert markdown to HTML. Does not need Gemini."""
    import markdown as md

    out_path = out_dir / f"{blob.hash}.html"

    title = analysis.get("title", "Document") if analysis else "Document"
    body_html = md.markdown(blob.text, extensions=["tables", "fenced_code", "toc"])

    html = f"""<!DOCTYPE html>
<html lang="{'zh' if has_cjk(blob.text) else 'en'}">
<head>
<meta charset="utf-8">
<title>{title}</title>
<style>
body {{ font-family: -apple-system, 'PingFang SC', sans-serif; max-width: 800px; margin: 2em auto; padding: 0 1em; line-height: 1.6; }}
table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
th {{ background: #f5f5f5; }}
code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; }}
pre {{ background: #f4f4f4; padding: 1em; overflow-x: auto; border-radius: 5px; }}
blockquote {{ border-left: 4px solid #ddd; margin: 1em 0; padding: 0.5em 1em; color: #666; }}
</style>
</head>
<body>
{body_html}
</body>
</html>"""

    out_path.write_text(html, encoding="utf-8")
    return out_path


def generate_pdf(blob: BlobInfo, analysis: dict, out_dir: Path) -> Optional[Path]:
    """Generate PDF from pdf_text (condensed executive summary)."""
    from fpdf import FPDF

    out_path = out_dir / f"{blob.hash}.pdf"
    pdf_text = analysis.get("pdf_text", blob.text[:3000])
    title = analysis.get("title", "Document")
    lang = analysis.get("language", "en")
    is_cjk = lang == "zh" or has_cjk(pdf_text)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    font_loaded = False
    if CJK_FONT_PATH:
        try:
            pdf.add_font("UniFont", "", CJK_FONT_PATH)
            pdf.set_font("UniFont", size=11)
            font_loaded = True
        except Exception:
            pass

    if not font_loaded:
        if is_cjk:
            return None
        pdf.set_font("Helvetica", size=11)
        pdf_text = pdf_text.encode("latin-1", errors="replace").decode("latin-1")
        title = title.encode("latin-1", errors="replace").decode("latin-1")

    # Title
    pdf.set_font_size(16)
    pdf.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font_size(11)

    # Body — condensed executive summary
    pdf.multi_cell(0, 6, pdf_text)

    pdf.output(str(out_path))
    return out_path


def generate_docx(blob: BlobInfo, analysis: dict, out_dir: Path) -> Optional[Path]:
    """Generate DOCX from docx_text (full detailed working document)."""
    from docx import Document
    from docx.shared import Pt

    out_path = out_dir / f"{blob.hash}.docx"
    title = analysis.get("title", "Document")
    docx_text = analysis.get("docx_text", blob.text[:10000])

    doc = Document()
    doc.add_heading(title, level=0)

    # Split docx_text into paragraphs by double newline
    paragraphs = [p.strip() for p in docx_text.split("\n\n") if p.strip()]
    for para_text in paragraphs:
        para = doc.add_paragraph(para_text)
        for run in para.runs:
            run.font.size = Pt(11)

    doc.save(str(out_path))
    return out_path


def generate_xlsx(analysis: dict, out_dir: Path, hash_id: str) -> Optional[Path]:
    """Generate XLSX from table_headers/table_rows if present."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    if not analysis.get("has_table"):
        return None

    headers = analysis.get("table_headers", [])
    rows = analysis.get("table_rows", [])
    if not headers or not rows:
        return None

    out_path = out_dir / f"{hash_id}.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = analysis.get("title", "Data")[:31]

    # Headers
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")

    # Data rows
    for row_idx, row_data in enumerate(rows, 2):
        for col_idx, val in enumerate(row_data, 1):
            if col_idx > len(headers):
                break
            # Try numeric conversion
            try:
                num = float(val)
                if num == int(num):
                    num = int(num)
                ws.cell(row=row_idx, column=col_idx, value=num)
            except (ValueError, TypeError):
                ws.cell(row=row_idx, column=col_idx, value=val)

    # Auto-width
    for col in ws.columns:
        max_len = 0
        col_letter = col[0].column_letter
        for cell in col:
            try:
                max_len = max(max_len, len(str(cell.value or "")))
            except Exception:
                pass
        ws.column_dimensions[col_letter].width = min(max_len + 2, 50)

    wb.save(str(out_path))
    return out_path


def generate_pptx(analysis: dict, out_dir: Path, hash_id: str) -> Optional[Path]:
    """Generate PPTX from pptx_slides (structured slide data)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    slides_data = analysis.get("pptx_slides", [])
    if len(slides_data) < 2:
        return None

    out_path = out_dir / f"{hash_id}.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = analysis.get("title", "Presentation")
    if slide.placeholders[1]:
        slide.placeholders[1].text = ""

    # Content slides from pptx_slides
    for s in slides_data[:15]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("slide_title", "")[:100]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = s.get("bullets", "")[:800]
        for para in tf.paragraphs:
            para.font.size = Pt(14)

    prs.save(str(out_path))
    return out_path


def generate_mp3(analysis: dict, out_dir: Path, hash_id: str) -> Optional[Path]:
    """Generate MP3 from mp3_script (conversational narration) via gTTS."""
    from gtts import gTTS

    script = analysis.get("mp3_script", "")
    if not script or len(script) < 20:
        return None

    out_path = out_dir / f"{hash_id}.mp3"
    lang = "zh-CN" if analysis.get("language") == "zh" else "en"

    try:
        tts = gTTS(text=script, lang=lang, slow=False)
        tts.save(str(out_path))
        return out_path
    except Exception as e:
        print(f"  [WARN] gTTS failed for {hash_id}: {e}")
        return None


def generate_png(analysis: dict, out_dir: Path, hash_id: str) -> Optional[Path]:
    """Generate PNG chart from chart data or quote image from key_quote."""
    key_quote = analysis.get("key_quote", "")
    out_path = out_dir / f"{hash_id}.png"
    lang = analysis.get("language", "en")

    if analysis.get("has_chart"):
        chart_data = {
            "type": analysis.get("chart_type", "bar"),
            "labels": analysis.get("chart_labels", []),
            "values": analysis.get("chart_values", []),
            "title": analysis.get("chart_title", "Chart"),
        }
        if chart_data["labels"] and chart_data["values"]:
            return _generate_chart_png(chart_data, out_path, lang)

    if key_quote and len(key_quote) > 15:
        return _generate_quote_png(key_quote, out_path, lang)
    return None


def _generate_chart_png(chart_data: dict, out_path: Path, lang: str) -> Optional[Path]:
    """Generate a chart PNG using matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    if lang == "zh":
        plt.rcParams["font.sans-serif"] = ["Songti SC", "STHeiti", "PingFang SC", "SimHei"]
        plt.rcParams["axes.unicode_minus"] = False

    chart_type = chart_data.get("type", "bar")
    labels = chart_data.get("labels", [])
    values = chart_data.get("values", [])
    title = chart_data.get("title", "Chart")

    if not labels or not values or len(labels) != len(values):
        return None

    fig, ax = plt.subplots(figsize=(8, 5))

    if chart_type == "pie":
        ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90)
    elif chart_type == "line":
        ax.plot(labels, values, marker="o", linewidth=2)
        ax.set_ylabel("Value")
        plt.xticks(rotation=45, ha="right")
    else:  # bar
        colors = plt.cm.Set2(range(len(labels)))
        ax.bar(labels, values, color=colors)
        ax.set_ylabel("Value")
        plt.xticks(rotation=45, ha="right")

    ax.set_title(title)
    plt.tight_layout()
    fig.savefig(str(out_path), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return out_path


def _generate_quote_png(quote: str, out_path: Path, lang: str) -> Optional[Path]:
    """Generate a quote card PNG using Pillow."""
    from PIL import Image, ImageDraw, ImageFont

    width, height = 800, 400
    bg_color = (45, 55, 72)
    text_color = (255, 255, 255)

    img = Image.new("RGB", (width, height), bg_color)
    draw = ImageDraw.Draw(img)

    # Try to find a suitable font
    font_size = 24
    font = None
    if lang == "zh" and CJK_FONT_PATH:
        try:
            font = ImageFont.truetype(CJK_FONT_PATH, font_size)
        except Exception:
            pass
    if font is None:
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
        except Exception:
            font = ImageFont.load_default()

    # Word wrap
    max_chars_per_line = 30 if lang == "zh" else 45
    lines = []
    for i in range(0, len(quote), max_chars_per_line):
        lines.append(quote[i:i + max_chars_per_line])
    lines = lines[:8]  # Cap lines

    # Draw centered
    total_height = len(lines) * (font_size + 10)
    y_start = (height - total_height) // 2

    # Decorative quote mark
    draw.text((30, y_start - 50), "\u201c", fill=(100, 140, 200), font=font)

    for i, line in enumerate(lines):
        bbox = draw.textbbox((0, 0), line, font=font)
        text_w = bbox[2] - bbox[0]
        x = (width - text_w) // 2
        y = y_start + i * (font_size + 10)
        draw.text((x, y), line, fill=text_color, font=font)

    img.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# Ground truth helpers
# ---------------------------------------------------------------------------

def _write_gt(gt_dir: Path, hash_id: str, modality: str, text: str):
    """Write ground truth text for a multimodal file."""
    gt_path = gt_dir / f"{hash_id}_{modality}_gt.blob"
    gt_path.write_text(text, encoding="utf-8")


def _format_xlsx_gt(analysis: dict) -> str:
    """Format table_headers + table_rows as TSV text for XLSX ground truth."""
    headers = analysis.get("table_headers", [])
    rows = analysis.get("table_rows", [])
    lines = ["\t".join(headers)]
    for row in rows:
        lines.append("\t".join(str(v) for v in row))
    return "\n".join(lines)


def _format_pptx_gt(analysis: dict) -> str:
    """Format pptx_slides as text for PPTX ground truth."""
    slides = analysis.get("pptx_slides", [])
    parts = []
    for s in slides:
        parts.append(f"[Slide] {s.get('slide_title', '')}\n{s.get('bullets', '')}")
    return "\n\n".join(parts)


def _format_png_gt(analysis: dict) -> str:
    """Format chart description or key quote for PNG ground truth."""
    if analysis.get("has_chart"):
        title = analysis.get("chart_title", "")
        labels = analysis.get("chart_labels", [])
        values = analysis.get("chart_values", [])
        return f"[Chart] {title}: {labels} = {values}"
    return f"[Quote] {analysis.get('key_quote', '')}"


# ---------------------------------------------------------------------------
# Core conversion logic
# ---------------------------------------------------------------------------

def convert_single_blob(client, blob: BlobInfo, out_dir: Path, gt_dir: Path) -> ConversionResult:
    """Convert a single blob to multiple modalities and write GT annotations."""
    result = ConversionResult(hash=blob.hash, original_path=blob.original_path)

    # 1. HTML always (no Gemini needed)
    try:
        p = generate_html(blob, None, out_dir)
        if p:
            result.modalities["html"] = {"path": str(p.name), "size": p.stat().st_size}
            _write_gt(gt_dir, blob.hash, "html", blob.text)
    except Exception as e:
        result.errors.append(f"html: {e}")

    # 2. Call Gemini for structured analysis
    analysis = None
    if client:
        analysis = analyze_blob(client, blob)

    if not analysis:
        # Gemini failed — only HTML was generated
        return result

    # Update HTML with analysis title
    try:
        html_path = out_dir / f"{blob.hash}.html"
        if html_path.exists():
            p = generate_html(blob, analysis, out_dir)
            if p:
                result.modalities["html"] = {"path": str(p.name), "size": p.stat().st_size}
    except Exception:
        pass

    # Gemini picks modalities; cap based on source text length
    # Short docs → 1, medium → up to 2, long/rich → up to 3
    text_len = len(blob.text)
    if text_len < 500:
        max_modalities = 1
    elif text_len < 2000:
        max_modalities = 2
    else:
        max_modalities = 3
    requested = set(list(analysis.get("modalities", []))[:max_modalities])

    # Map modality -> (generator, gt_text)
    gt_texts = {
        "pdf":  analysis.get("pdf_text", ""),
        "docx": analysis.get("docx_text", ""),
        "xlsx": _format_xlsx_gt(analysis) if analysis.get("has_table") else "",
        "pptx": _format_pptx_gt(analysis),
        "mp3":  analysis.get("mp3_script", ""),
        "png":  _format_png_gt(analysis),
    }

    generators = {
        "pdf":  lambda: generate_pdf(blob, analysis, out_dir),
        "docx": lambda: generate_docx(blob, analysis, out_dir),
        "xlsx": lambda: generate_xlsx(analysis, out_dir, blob.hash),
        "pptx": lambda: generate_pptx(analysis, out_dir, blob.hash),
        "mp3":  lambda: generate_mp3(analysis, out_dir, blob.hash),
        "png":  lambda: generate_png(analysis, out_dir, blob.hash),
    }

    for modality, gen_fn in generators.items():
        if modality not in requested:
            continue
        try:
            p = gen_fn()
            if p and p.exists():
                result.modalities[modality] = {"path": str(p.name), "size": p.stat().st_size}
                # Write GT annotation for this modality
                gt_text = gt_texts.get(modality, "")
                if gt_text:
                    _write_gt(gt_dir, blob.hash, modality, gt_text)
        except Exception as e:
            result.errors.append(f"{modality}: {e}")
            if os.environ.get("DEBUG"):
                traceback.print_exc()

    return result


def _progress(traj_idx: int, traj_total: int, traj_name: str,
              blob_done: int, blob_total: int, files_so_far: int):
    """Single-line progress overwrite."""
    sys.stdout.write(
        f"\r[{traj_idx}/{traj_total}] {traj_name} | "
        f"blob {blob_done}/{blob_total} | "
        f"{files_so_far} files total"
        "        "  # padding to clear previous longer lines
    )
    sys.stdout.flush()


def convert_trajectory(traj_name: str, client, parallel: int = 4,
                       skip_existing: bool = False, dry: bool = False,
                       traj_idx: int = 0, traj_total: int = 0,
                       files_so_far: int = 0) -> dict:
    """Convert all blobs in one trajectory."""
    traj_dir = SIGNAL_DIR / traj_name
    if not traj_dir.exists():
        print(f"[SKIP] {traj_name}: directory not found")
        return {"trajectory": traj_name, "status": "not_found"}

    blobs = load_content_blobs(traj_dir)
    if not blobs:
        print(f"[SKIP] {traj_name}: no content blobs")
        return {"trajectory": traj_name, "status": "no_blobs"}

    out_dir = traj_dir / "media" / "blobs_multimodal"
    gt_dir = traj_dir / "media" / "blobs_multimodal_gt"

    if dry:
        print(f"[DRY] {traj_name}: {len(blobs)} blobs -> {out_dir}")
        return {"trajectory": traj_name, "status": "dry", "blob_count": len(blobs)}

    out_dir.mkdir(parents=True, exist_ok=True)
    gt_dir.mkdir(parents=True, exist_ok=True)

    # Check existing manifest
    manifest_path = out_dir / "manifest.json"
    existing_hashes = set()
    if skip_existing and manifest_path.exists():
        try:
            with open(manifest_path) as f:
                existing = json.load(f)
            existing_hashes = set(existing.get("entries", {}).keys())
        except Exception:
            pass

    # Filter out already-converted
    if skip_existing and existing_hashes:
        blobs = [b for b in blobs if b.hash not in existing_hashes]
        if not blobs:
            print(f"[SKIP] {traj_name}: all {len(existing_hashes)} blobs already converted")
            return {"trajectory": traj_name, "status": "already_done"}

    t0 = time.time()

    results = []
    blob_done = 0
    blob_total = len(blobs)
    blob_timeout = 90  # seconds per blob

    _progress(traj_idx, traj_total, traj_name, 0, blob_total, files_so_far)

    if parallel <= 1:
        for blob in blobs:
            try:
                r = convert_single_blob(client, blob, out_dir, gt_dir)
                results.append(r)
            except Exception as e:
                results.append(ConversionResult(
                    hash=blob.hash, original_path=blob.original_path,
                    errors=[str(e)],
                ))
            blob_done += 1
            _progress(traj_idx, traj_total, traj_name, blob_done, blob_total, files_so_far)
    else:
        with ThreadPoolExecutor(max_workers=parallel) as executor:
            futures = {
                executor.submit(convert_single_blob, client, blob, out_dir, gt_dir): blob
                for blob in blobs
            }
            for future in as_completed(futures):
                blob = futures[future]
                try:
                    r = future.result(timeout=blob_timeout)
                    results.append(r)
                except TimeoutError:
                    results.append(ConversionResult(
                        hash=blob.hash, original_path=blob.original_path,
                        errors=["timeout"],
                    ))
                except Exception as e:
                    results.append(ConversionResult(
                        hash=blob.hash, original_path=blob.original_path,
                        errors=[str(e)],
                    ))
                blob_done += 1
                _progress(traj_idx, traj_total, traj_name, blob_done, blob_total, files_so_far)

    elapsed = time.time() - t0

    # Build manifest
    manifest_entries = {}
    # Load existing entries first
    if skip_existing and manifest_path.exists():
        try:
            with open(manifest_path) as f:
                old = json.load(f)
            manifest_entries = old.get("entries", {})
        except Exception:
            pass

    total_files = 0
    total_errors = 0
    modality_counts = {}

    for r in results:
        if not r.modalities:
            # Skip failed/timed-out blobs so -s will retry them
            total_errors += len(r.errors)
            continue
        entry = {
            "original_path": r.original_path,
            "modalities": r.modalities,
        }
        if r.errors:
            entry["errors"] = r.errors
            total_errors += len(r.errors)
        manifest_entries[r.hash] = entry
        total_files += len(r.modalities)
        for m in r.modalities:
            modality_counts[m] = modality_counts.get(m, 0) + 1

    manifest = {
        "version": 1,
        "trajectory": traj_name,
        "entries": manifest_entries,
        "stats": {
            "total_blobs_processed": len(results),
            "total_files_generated": total_files,
            "total_errors": total_errors,
            "modality_counts": modality_counts,
            "elapsed_seconds": round(elapsed, 1),
        },
    }

    with open(manifest_path, "w") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)

    print(f"  Done: {total_files} files in {elapsed:.1f}s | {modality_counts}")
    return manifest


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def discover_trajectories() -> list[str]:
    """Find all trajectory directories in signal/."""
    if not SIGNAL_DIR.exists():
        return []
    return sorted([
        d.name for d in SIGNAL_DIR.iterdir()
        if d.is_dir() and not d.name.startswith("_") and not d.name.startswith(".")
        and (d / "media" / "manifest.json").exists()
    ])


def main():
    parser = argparse.ArgumentParser(description="Convert text blobs to multimodal files")
    parser.add_argument("--trajectory", "-t", help="Single trajectory name to process")
    parser.add_argument("--parallel", "-p", type=int, default=4, help="Parallel workers per trajectory (default: 4)")
    parser.add_argument("--skip-existing", "-s", action="store_true", help="Skip already-converted blobs")
    parser.add_argument("--dry", action="store_true", help="Dry run: show what would be converted")
    args = parser.parse_args()

    # Discover trajectories
    if args.trajectory:
        trajectories = [args.trajectory]
    else:
        trajectories = discover_trajectories()
        # Filter out non-profile dirs
        trajectories = [t for t in trajectories if re.match(r"p\d+_", t)]

    if not trajectories:
        print("No trajectories found.")
        return

    print(f"Found {len(trajectories)} trajectories")

    # Create Gemini client (unless dry run)
    client = None
    if not args.dry:
        try:
            client = create_gemini_client()
            print("Gemini client initialized")
        except Exception as e:
            print(f"[WARN] Gemini unavailable: {e}")
            print("  Will only generate HTML (no Gemini-dependent modalities)")

    # Process trajectories sequentially (avoid Gemini rate limits)
    summary = {"total_trajectories": 0, "total_files": 0, "total_errors": 0}
    t_global = time.time()

    n_traj = len(trajectories)
    for i, traj in enumerate(trajectories, 1):
        result = convert_trajectory(
            traj, client,
            parallel=args.parallel,
            skip_existing=args.skip_existing,
            dry=args.dry,
            traj_idx=i,
            traj_total=n_traj,
            files_so_far=summary["total_files"],
        )
        summary["total_trajectories"] += 1
        stats = result.get("stats", {})
        summary["total_files"] += stats.get("total_files_generated", 0)
        summary["total_errors"] += stats.get("total_errors", 0)
    print()  # newline after \r progress

    elapsed_total = time.time() - t_global
    print(f"\n{'='*60}")
    print(f"DONE: {summary['total_trajectories']} trajectories, "
          f"{summary['total_files']} files, "
          f"{summary['total_errors']} errors, "
          f"{elapsed_total:.0f}s total")


if __name__ == "__main__":
    main()
